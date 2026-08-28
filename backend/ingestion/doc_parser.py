"""
Yonder Graph — Multi-Format Document Parser

Extracts clean text and structural content from various file formats:
  - PDF (.pdf)
  - Presentations (.pptx, .ppt)
  - Spreadsheets (.xlsx, .xls, .csv)
  - Word documents (.docx, .doc)
  - Markdown & Text (.md, .txt)
"""

import io
import logging
from typing import Tuple

logger = logging.getLogger(__name__)


def extract_text_from_file(filename: str, file_bytes: bytes) -> Tuple[str, str]:
    """
    Extract text from file bytes based on file extension.
    
    Returns:
        (extracted_text, file_type)
    """
    ext = "." + filename.split(".")[-1].lower() if "." in filename else ".txt"
    text = ""

    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            pages = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages.append(f"--- Page {i+1} ---\n{page_text}")
            text = "\n\n".join(pages)

        elif ext in (".docx", ".doc"):
            try:
                import docx
                doc_file = io.BytesIO(file_bytes)
                doc = docx.Document(doc_file)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                # Also extract table text
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            paragraphs.append(row_text)
                text = "\n".join(paragraphs)
            except Exception as e:
                # Fallback for binary .doc or partial docx
                logger.warning("docx parser failed, falling back to raw decode: %s", e)
                text = file_bytes.decode("utf-8", errors="ignore")

        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                ppt_file = io.BytesIO(file_bytes)
                prs = Presentation(ppt_file)
                slides = []
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_texts.append(shape.text.strip())
                    if slide_texts:
                        slides.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))
                text = "\n\n".join(slides)
            except Exception as e:
                logger.warning("pptx parser failed, falling back: %s", e)
                text = file_bytes.decode("utf-8", errors="ignore")

        elif ext in (".xlsx", ".xls"):
            try:
                import pandas as pd
                excel_file = io.BytesIO(file_bytes)
                excel = pd.read_excel(excel_file, sheet_name=None)
                sheets = []
                for sheet_name, df in excel.items():
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + df.to_string(index=False))
                text = "\n\n".join(sheets)
            except Exception as e:
                logger.warning("excel parser failed: %s", e)
                text = file_bytes.decode("utf-8", errors="ignore")

        elif ext == ".csv":
            try:
                import pandas as pd
                csv_file = io.BytesIO(file_bytes)
                df = pd.read_csv(csv_file)
                text = df.to_string(index=False)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")

        else:
            # Plain text, markdown, json, etc.
            text = file_bytes.decode("utf-8", errors="ignore")

    except Exception as e:
        logger.error("Failed to parse file %s: %s", filename, e)
        text = file_bytes.decode("utf-8", errors="ignore")

    return text.strip(), ext
